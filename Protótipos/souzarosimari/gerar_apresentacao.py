from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Cria a pasta de imagens se não existir
os.makedirs('static', exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)

def add_image_slide(title, image_path, desc=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Layout em branco
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(3), Inches(1.5), height=Inches(6))
    else:
        txt = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(10), Inches(1))
        txt.text_frame.text = f"Imagem não encontrada: {image_path}\nSalve as imagens geradas na pasta /static"
    
    if desc:
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(7.8), Inches(14), Inches(0.8))
        desc_box.text_frame.text = desc
        desc_box.text_frame.paragraphs[0].font.size = Pt(18)

# SLIDES DA APRESENTAÇÃO
add_title_slide("Sistema de Catraca Virtual para Ônibus", 
                "Demonstração com Python, Flask e SQLite\nGitHub Codespaces | Outubro 2026")

add_content_slide("1. Objetivo do Projeto", [
    "Simular catraca de ônibus 100% web, sem hardware físico",
    "Permitir criação de cartões virtuais e validação por ID",
    "Demonstrar: banco de dados, validação, múltiplos dispositivos",
    "Uso acadêmico para conceitos de sistemas distribuídos"
])

add_content_slide("2. Tecnologias Utilizadas", [
    "Back-end: Python 3 + Flask",
    "Front-end: HTML5 + CSS3 Responsivo Mobile-First",
    "Banco de Dados: SQLite com 2 tabelas",
    "Ambiente: GitHub Codespaces com porta pública",
    "Demonstração: 2 celulares simultâneos"
])

add_content_slide("3. MVP Entregue - Requisitos Funcionais", [
    "RF01: Criar cartão com nome e saldo inicial",
    "RF02: Gerar ID único automático para cada cartão",
    "RF03: Consultar cartão por ID",
    "RF04: Catraca virtual valida entrada",
    "RF05: Aprova se cartão ativo e saldo >= R$ 5,00",
    "RF06: Nega se inexistente, bloqueado ou sem saldo",
    "RF07: Desconta R$ 5,00 se aprovado",
    "RF08: Registra todas as tentativas no histórico"
])

add_content_slide("4. Regras de Negócio", [
    "1. Cartão só libera se status = 'Ativo'",
    "2. Saldo menor que R$ 5,00 = Entrada Negada",
    "3. ID inexistente = Entrada Negada",
    "4. Valor da passagem fixo: R$ 5,00",
    "5. Toda tentativa é registrada, mesmo negada"
])

# SLIDES COM IMAGENS DAS TELAS
add_image_slide("5. Tela 1: Página Inicial", "static/tela1_inicio.png", 
                "Menu principal com acesso para Passageiro e Operador da Catraca")

add_image_slide("6. Tela 2: Criar Cartão Virtual", "static/tela2_criar.png",
                "Formulário: nome do passageiro e saldo inicial")

add_image_slide("7. Tela 3: Cartão Criado", "static/tela3_cartao.png",
                "Sistema exibe ID 1001, saldo R$ 20,00 e status Ativo")

add_image_slide("8. Tela 4: Catraca Virtual", "static/tela4_catraca.png",
                "Operador digita o ID do cartão para validar")

add_image_slide("9. Tela 5: Entrada Aprovada", "static/tela5_aprovada.png",
                "Saldo suficiente: desconta R$ 5,00 e libera entrada")

add_image_slide("10. Tela 6: Entrada Negada", "static/tela6_negada.png",
                "Motivo: Saldo insuficiente, cartão bloqueado ou inexistente")

add_image_slide("11. Tela 7: Histórico de Passagens", "static/tela7_historico.png",
                "Tabela com data, ID, status, motivo e valor cobrado")

add_content_slide("12. Roteiro de Demonstração", [
    "1. PC: Abrir Codespaces e tornar porta 5000 pública",
    "2. Celular 1: Criar cartão 'João Silva' com R$ 20,00 → ID 1001",
    "3. Celular 2: Catraca → digitar 1001 → Validar",
    "4. Mostrar: 'Entrada Aprovada. Boa viagem!'",
    "5. Celular 1: Consultar cartão → Saldo agora R$ 15,00",
    "6. Mostrar histórico + testar ID 9999 = Negado"
])

add_content_slide("13. Conclusão", [
    "MVP 100% funcional entregue conforme PRD",
    "Simula fluxo completo: cadastro → validação → registro",
    "Riscos mitigados: CSS responsivo + init_db() + teste de porta",
    "Próximos passos: QR Code, tela admin, login, relatórios"
])

prs.save('Apresentacao_Catraca_Virtual.pptx')
print("✅ PowerPoint gerado: Apresentacao_Catraca_Virtual.pptx")
